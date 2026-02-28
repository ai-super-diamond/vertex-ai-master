#!/usr/bin/env python3
"""Generate PowerPoint presentation for Vertex AI Master CLI Architecture"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Generate comprehensive architecture presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Vertex AI Master CLI"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 102, 204)
    title_p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Project Architecture Overview"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Java 25 • Picocli • Google Cloud Vertex AI"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(16)
    footer_p.font.italic = True
    footer_p.font.color.rgb = RGBColor(150, 150, 150)
    footer_p.alignment = PP_ALIGN.CENTER

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Overview"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Vertex AI Master CLI"

    p = tf.add_paragraph()
    p.text = "Command-line interface for Google's Vertex AI generative models"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Built with Java 25, Picocli, and clean layered architecture"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Supports dual API routing (Vertex AI SDK + Chat Completions API)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Tests model availability across 40+ GCP regions"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Native executable support via GraalVM"
    p.level = 1

    # Slide 3: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Dual API Support"

    p = tf.add_paragraph()
    p.text = "Seamlessly switches between Vertex AI SDK and Chat Completions API"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Region Availability Testing"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Check model availability across clusters (US, EU, ASIA, etc.)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Model Alias System"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Define short names for complex model IDs in models.properties"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Flexible Authentication"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "API Key, Service Account JSON, Application Default Credentials"
    p.level = 1

    # Slide 4: 3-Tier Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "3-Tier Layered Architecture"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Presentation Layer (CLI)"

    p = tf.add_paragraph()
    p.text = "VertexAiMasterMain - Picocli-based CLI"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Handles user input, authentication options, region checks"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "2. Service Layer"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "VertexAiService - Business logic and orchestration"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Model resolution, region management, request routing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "3. Client Layer"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "VertexAiClient - Direct API communication"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ChatCompletionsClient, WorldwideAvailabilityClient"
    p.level = 1

    # Slide 5: Core Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Core Components"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "VertexAiClient"

    p = tf.add_paragraph()
    p.text = "Main client for API communication"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Automatic routing: Standard API vs Chat Completions"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ChatCompletionsClient"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Handles Model-as-a-Service (MaaS) offerings (DeepSeek, Qwen, etc.)"
    p.level = 1


    p = tf.add_paragraph()
    p.text = "WorldwideAvailabilityClient"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Tests model availability across all global regions"
    p.level = 1

    # Slide 6: Data Transfer Objects (DTOs)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Transfer Objects (DTOs)"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AuthenticationConfig"

    p = tf.add_paragraph()
    p.text = "API Key, Service Account, ADC configurations"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "GenerationRequest / GenerationResult"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Request/response for content generation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "RegionCheckRequest / RegionCheckResult"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Region availability testing data structures"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "ErrorType enum"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Categorizes API errors for better diagnostics"
    p.level = 1

    # Slide 7: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Language & Runtime"

    p = tf.add_paragraph()
    p.text = "Java 25 (latest LTS features)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "CLI Framework"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Picocli 4.7.7 - Annotation-based command parsing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Google Cloud SDKs"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Google GenAI SDK 1.32.0, Google Auth Library 1.41.0"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Build & Tooling"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Maven 3.9+, GraalVM (native builds), Spotless (formatting)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Testing"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "JUnit 5, AssertJ, Mockito"
    p.level = 1

    # Slide 8: Authentication Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Authentication Flow"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Three Authentication Methods"

    p = tf.add_paragraph()
    p.text = "1. API Key (Gemini API)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Direct API key for public Gemini models"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "No GCP project required"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "2. Service Account with Explicit Key File"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Loads credentials from JSON key file"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Full Vertex AI capabilities"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "3. Application Default Credentials (ADC)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Uses GOOGLE_APPLICATION_CREDENTIALS environment variable"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Seamless local development and cloud deployment"
    p.level = 1

    # Slide 9: Model Routing Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Model Routing Strategy"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Intelligent API Selection"

    p = tf.add_paragraph()
    p.text = "Standard Vertex AI API"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Gemini, Llama, Claude models"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Uses GenAI SDK with generateContent()"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Chat Completions API (MaaS)"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "DeepSeek, Qwen models"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Configured via .provider property in models.properties"
    p.level = 1


    # Slide 10: Region Management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Region Management"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "RegionCatalog & RegionProvider"

    p = tf.add_paragraph()
    p.text = "Supports 40+ GCP regions organized by clusters"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Clusters: US, EU, ASIA, MIDDLE_EAST, AFRICA, CANADA, SOUTH_AMERICA"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Cluster-Wide Testing"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "--check-all-regions --cluster US"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Worldwide Testing"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "--worldwide flag tests all regions globally"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Per-Model Region Override"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "models.properties supports .region=global for global endpoints"
    p.level = 1

    # Slide 11: Build System & Native Compilation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Build System & Native Compilation"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Maven Multi-Profile Build"

    p = tf.add_paragraph()
    p.text = "Default: JAR with dependencies (maven-shade-plugin)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Profile 'native': GraalVM native executable (vertex.exe)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Profile 'spotbugs': Static analysis"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Code Quality Tools"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Spotless: Auto-formatting with Eclipse formatter"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Checkstyle: Google Java Style enforcement"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "OpenRewrite: Dependency version management"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Native Build Benefits"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Near-instant startup, single executable, no JVM required"
    p.level = 1

    # Slide 12: Testing Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Testing Architecture"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Comprehensive Test Suite"

    p = tf.add_paragraph()
    p.text = "Unit Tests"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "VertexAiClientTest, RegionCatalogTest, AuthenticationConfigTest"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Mockito with inline mocking for final classes"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Integration Tests"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "WorldwideAvailabilityClientTest"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "VertexAiMasterMainTest - CLI integration"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Testing Tools"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "JUnit 5 with AssertJ fluent assertions"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Byte Buddy agent for advanced mocking scenarios"
    p.level = 1

    # Slide 13: Utilities & Supporting Classes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Utilities & Supporting Classes"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "PropertiesLoader"

    p = tf.add_paragraph()
    p.text = "Loads models.properties with system property override support"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "MarkdownReportGenerator"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Generates formatted Markdown reports for region availability"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Includes model metadata, success/fail counts, region details"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "VertexUtils"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Common utility functions for API interaction"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Logging Infrastructure"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "SLF4J API with Logback Classic implementation"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Configured via logback.xml in resources"
    p.level = 1

    # Slide 14: Configuration Management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Configuration Management"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "models.properties Structure"

    p = tf.add_paragraph()
    p.text = "Model Aliases"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "gemini.pro=gemini-1.5-pro-001"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Provider Mapping"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "deepseek.r1.0528.provider=deepseek-ai"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Region Override"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "gemini.pro.region=us-central1"
    p.level = 1


    p = tf.add_paragraph()
    p.text = "External Configuration"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "System property: models.config=path/to/custom.properties"
    p.level = 1

    # Slide 15: Summary & Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary & Architecture Benefits"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Clean Separation of Concerns"

    p = tf.add_paragraph()
    p.text = "CLI ↔ Service ↔ Client layers with DTOs"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Extensibility"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Easy to add new models, providers, or API endpoints"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Flexibility"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Multiple auth methods, configurable model routing"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Testability"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Comprehensive unit and integration test coverage"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Performance"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Native compilation for instant startup and low resource usage"
    p.level = 1

    # Save presentation
    output_file = "Vertex_AI_Master_CLI_Architecture.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_presentation()
